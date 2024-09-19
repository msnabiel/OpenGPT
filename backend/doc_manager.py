import os
import shutil
import pdfplumber
import pytesseract
from PIL import Image
from pptx import Presentation

# Ensure the /documents directory exists
os.makedirs("documents", exist_ok=True)

def check_file_exists(txt_filename):
    """Check if the file already exists in the documents folder"""
    return os.path.exists(f"documents/{txt_filename}")

def get_unique_filename(original_filename):
    """Prompt the user for a new filename if the original one already exists"""
    while True:
        new_name = input(f"File '{original_filename}' already exists. Enter a new filename (without extension): ")
        new_txt_filename = new_name + ".txt"
        if not check_file_exists(new_txt_filename):
            return new_txt_filename
        else:
            print(f"File '{new_txt_filename}' also exists. Try again.")

def copy_txt_file(txt_path, txt_filename):
    """Copy a .txt file to the documents directory"""
    if check_file_exists(txt_filename):
        txt_filename = get_unique_filename(txt_filename)
    
    shutil.copy(txt_path, f"documents/{txt_filename}")
    print(f"[INFO] Copied {txt_path} to /documents as {txt_filename}")

def extract_text_from_pdf(pdf_path, txt_filename):
    """Extract text from a PDF file and save as .txt"""
    print(f"[INFO] Extracting text from PDF: {pdf_path}")
    with pdfplumber.open(pdf_path) as pdf:
        text = ""
        for page in pdf.pages:
            text += page.extract_text()
    
    if check_file_exists(txt_filename):
        txt_filename = get_unique_filename(txt_filename)
    
    with open(f"documents/{txt_filename}", "w") as txt_file:
        txt_file.write(text)
    
    print(f"[INFO] PDF extraction completed. Saved as {txt_filename}")

def extract_text_from_ppt(ppt_path, txt_filename):
    """Extract text from a PPT file and save as .txt"""
    print(f"[INFO] Extracting text from PPT: {ppt_path}")
    prs = Presentation(ppt_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    
    if check_file_exists(txt_filename):
        txt_filename = get_unique_filename(txt_filename)
    
    with open(f"documents/{txt_filename}", "w") as txt_file:
        txt_file.write(text)

    print(f"[INFO] PPT extraction completed. Saved as {txt_filename}")

def extract_text_from_image(image_path, txt_filename):
    """Extract text from an image using OCR and save as .txt"""
    print(f"[INFO] Extracting text from image: {image_path}")
    image = Image.open(image_path)
    text = pytesseract.image_to_string(image)
    
    if check_file_exists(txt_filename):
        txt_filename = get_unique_filename(txt_filename)
    
    with open(f"documents/{txt_filename}", "w") as txt_file:
        txt_file.write(text)

    print(f"[INFO] Image text extraction completed. Saved as {txt_filename}")

def convert_and_save(file_path):
    """Convert a file (PDF, PPT, image) to a .txt file and save"""
    filename, extension = os.path.splitext(file_path)
    base_filename = os.path.basename(filename)
    txt_filename = base_filename + ".txt"

    print(f"[INFO] Processing file: {file_path}")
    
    if extension.lower() == ".pdf":
        extract_text_from_pdf(file_path, txt_filename)
    elif extension.lower() in [".pptx", ".ppt"]:
        extract_text_from_ppt(file_path, txt_filename)
    elif extension.lower() in [".jpg", ".jpeg", ".png", ".tiff", ".bmp"]:
        extract_text_from_image(file_path, txt_filename)
    elif extension.lower() == ".txt":
        copy_txt_file(file_path, txt_filename)  # Just copy .txt files to /documents
    else:
        print(f"[WARNING] Unsupported file type: {extension}")
        return

    # Commented out section for deleting the original file
    # os.remove(file_path)
    print(f"[INFO] Processed and saved: {file_path} as {txt_filename}")

def process_files_in_directory(directory):
    """Process all files in a directory and convert them to .txt"""
    files = os.listdir(directory)
    total_files = len(files)
    print(f"[INFO] Found {total_files} files in {directory}")

    processed_count = 0
    for file in files:
        file_path = os.path.join(directory, file)
        if os.path.isfile(file_path):
            print(f"[INFO] Processing file {processed_count + 1}/{total_files}: {file}")
            convert_and_save(file_path)
            processed_count += 1
        else:
            print(f"[WARNING] {file} is not a file. Skipping...")

    print(f"[INFO] Processed {processed_count} files out of {total_files}")

# Set the directory where your files (PDF, PPT, images, txt) are located
input_directory = "database"
process_files_in_directory(input_directory)